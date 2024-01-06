"""
complex_assignment_14:

Write A Python Program To Create Powerpoint File On User Request on Web 

"""
from pptx import Presentation

def create_powerpoint_file(title, content):
    # Create a presentation object
    presentation = Presentation()

    # Add a slide with a title and content layout
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)

    # Set the title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Set the content
    content_placeholder = slide.placeholders[1]
    content_placeholder.text = content

    # Save the presentation to a file
    presentation.save('E:/bulding/python_exercises/3-Advanced/Study Matrials Complex Problem14/user_presentation.pptx')

# Get user input
title = input("Enter the title for the slide: ")
content = input("Enter the content for the slide: ")

# Create the PowerPoint file
create_powerpoint_file(title, content)